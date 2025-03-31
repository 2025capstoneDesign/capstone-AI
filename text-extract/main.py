import os
import whisper
import yt_dlp
from pptx import Presentation
from tqdm import tqdm  



def ensure_downloads_dir():
    if not os.path.exists(DOWNLOAD_DIR):
        os.makedirs(DOWNLOAD_DIR)

def download_youtube_audio(youtube_url):
    """
    yt-dlp를 사용하여 YouTube 영상에서 오디오 추출 
    """
    ensure_downloads_dir()
    audio_path = os.path.join(DOWNLOAD_DIR, "audio.wav")
    
    print("\n📥 유튜브 영상 오디오 다운로드 중...")
    ydl_opts = {
        "format": "bestaudio/best",
        "outtmpl": audio_path,
    }

    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        ydl.download([youtube_url])
    
    print("\n✅ 유튜브 영상 오디오 다운로드 완료!")
    return audio_path

def transcribe_audio(audio_path, model_size="small"):
    """
    오디오를 텍스트로 변환
    """
    print("\n 오디오를 텍스트로 변환 중...")
    model = whisper.load_model(model_size)

    for _ in tqdm(range(100), desc="오디오 텍스트 변환 진행 중", unit="%"):
        pass

    result = model.transcribe(audio_path)
    print("\n 오디오를 텍스트로 변환 완료!")
    return result["text"]

def extract_ppt_text(ppt_path):
    """
    주어진 PPT 파일에서 슬라이드별 텍스트를 추출하는 함수
    """
    print("\n PPT에서 텍스트 추출 중...")
    prs = Presentation(ppt_path)
    slide_texts = []

    for i, slide in enumerate(tqdm(prs.slides, desc="슬라이드 분석 진행 중", unit="slide")):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())

        slide_texts.append(f"슬라이드 {i+1}:\n" + "\n".join(slide_content) + "\n" + "-"*40)

    print("\n PPT 텍스트 추출 완료!")
    return slide_texts


# 다운로드 폴더 준비
DOWNLOAD_DIR = "download"
ensure_downloads_dir()

# PPT에서 텍스트 추출
ppt_file = "./test.pptx"
slide_contents = extract_ppt_text(ppt_file)
ppt_text_path = os.path.join(DOWNLOAD_DIR, "ppt_text.txt")

with open(ppt_text_path, "w", encoding="utf-8") as f:
    for slide in slide_contents:
        f.write(slide + "\n\n")

print(f"\n 추출한 PPT 텍스트를 '{ppt_text_path}'에 저장하였습니다.")

# 유튜브 오디오 다운로드 및 변환
youtube_url = "https://www.youtube.com/watch?v=74sEFYBBRAY&t=14s"
audio_path = download_youtube_audio(youtube_url)
transcribed_text = transcribe_audio(audio_path)

lecture_text_path = os.path.join(DOWNLOAD_DIR, "lecture_text.txt")
with open(lecture_text_path, "w", encoding="utf-8") as f:
    f.write(transcribed_text)

print(f"\n 추출한 강의 음성 텍스트를 '{lecture_text_path}'에 저장하였습니다.")
